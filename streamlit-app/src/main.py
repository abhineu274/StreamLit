import streamlit as st
from openai import AzureOpenAI
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
import requests
import base64
from io import BytesIO
import json

# Load environment variables
load_dotenv()

endpoint = os.getenv("ENDPOINT_URL", "https://aoai-d01.openai.azure.com/")
deployment = os.getenv("DEPLOYMENT_NAME", "gpt-4o")
image_deployment = os.getenv("DALL_E_DEPLOYMENT_NAME", "dall-e-3")

client = AzureOpenAI(
    azure_endpoint=endpoint,
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    api_version="2024-05-01-preview",
)

# Function to generate slide content
def generate_slide_content(topic, num_slides):
    chat_prompt = [
        {
            "role": "system",
            "content": f"Generate a presentation on {topic} with {num_slides} slides. Also create relevant DALL-E images. Return the slides in json format, so that json.loads() can be used to parse the response."
        }
    ]

    response = client.chat.completions.create(
        model=deployment,
        messages=chat_prompt,
        max_tokens=1000
    )
    slides = response.choices[0].message.content
    st.write("API Response:", slides)  # Debugging line
    if not slides:
        st.error("Received empty response from the API.")
        return []
    try:
        print('slides :', slides)
        slides_json = json.loads(slides)  # Convert JSON to Python list
        print('slides_json :', slides_json)
        return slides_json
    except json.JSONDecodeError as e:
        st.error(f"Failed to parse the response from the API: {e}")
        return []

# Function to generate images using DALL·E
def generate_images(slides, topic):
    print('Generating images for slides...')
    for slide in slides:
        print(f"Generating image for slide '{slide['title']}'...")
        image_prompt = f"An engaging, high-quality image related to: {slide['title']}."
        try:
            response = client.images.generate(
                model=image_deployment,
                prompt=image_prompt,
                n=1,
                size="1024x1024"
            )
            print('response :', response)
            if isinstance(slide, dict):
                slide["image_url"] = response.data[0].url  # Attach image URL
            else:
                st.error(f"Slide is not a dictionary: {slide}")
        except Exception as e:
            st.error(f"An unexpected error occurred while generating image for slide '{slide['title']}']: {e}")
            if isinstance(slide, dict):
                slide["image_url"] = None  # Set image URL to None if generation fails
    return slides

# Function to create PPT
def create_ppt(slides):
    prs = Presentation()
    for slide_data in slides:
        slide_layout = prs.slide_layouts[5]  # Title & Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = slide_data["title"]
        content.text = "\n".join(slide_data["content"])
        # Fetch DALL·E image
        img_url = slide_data["image_url"]
        if img_url:
            img_response = requests.get(img_url)
            img_stream = BytesIO(img_response.content)
            left = Inches(5)
            top = Inches(1)
            slide.shapes.add_picture(img_stream, left, top, width=Inches(3))
    # Save presentation
    pptx_file = "generated_presentation.pptx"
    prs.save(pptx_file)
    return pptx_file

# Streamlit UI
st.title("Azure OpenAI Multi-Agent PPT Generator")

with st.form(key="ppt_form"):
    topic = st.text_input("Enter PPT Topic")
    num_slides = st.number_input("Number of Slides", min_value=1, step=1)
    submit_button = st.form_submit_button(label="Generate Presentation")

if submit_button:
    if topic and num_slides:
        print("program started")
        slides = generate_slide_content(topic, num_slides)
        if slides:
            slides = generate_images(slides, topic)
            pptx_file = create_ppt(slides)
            # Provide download link
            with open(pptx_file, "rb") as f:
                bytes_data = f.read()
                b64 = base64.b64encode(bytes_data).decode()
                href = f'<a href="data:file/pptx;base64,{b64}" download="generated_presentation.pptx">Download Presentation</a>'
                st.markdown(href, unsafe_allow_html=True)
    else:
        st.warning("Please enter both topic and number of slides.")